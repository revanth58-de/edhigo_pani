import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Constants
BG_COLOR = RGBColor(249, 249, 249)
CARD_COLOR = RGBColor(255, 255, 255)
TEXT_DARK = RGBColor(33, 33, 33)
TEXT_MUTED = RGBColor(97, 97, 97)
TEXT_LIGHT = RGBColor(255, 255, 255)

BRAND_GREEN = RGBColor(46, 125, 50)     # AP/Telangana/Karnataka Agri Green
BRAND_BLUE = RGBColor(21, 101, 192)      # Tech Blue
BRAND_ORANGE = RGBColor(239, 108, 0)     # Mestri Orange
BRAND_RED = RGBColor(198, 40, 40)        # Pain point Red
LINE_COLOR = RGBColor(224, 224, 224)

# Paths
ARTIFACT_DIR = r"C:\Users\renan\.gemini\antigravity\brain\309ad51f-cdf5-4ba7-a9f0-3b541e6d7118"
LOGO_PATH = r"c:\Users\renan\OneDrive\Desktop\edhigo_pani\mobile\assets\logo_transparent.png"
OUTPUT_PATH = r"c:\Users\renan\OneDrive\Desktop\edhigo_pani\dinasari_pitch_deck_v2.pptx"

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, margin=0):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    return tf

def add_header_footer(slide, category, title):
    set_slide_background(slide, BG_COLOR)
    
    # Header Textbox
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10), Inches(1.0))
    tf = header_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    # Category
    p_cat = tf.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.name = "Segoe UI"
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = BRAND_GREEN
    
    # Title
    p_title = tf.add_paragraph()
    p_title.text = title
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_DARK
    p_title.space_before = Pt(2)
    
    # Logo
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(11.8), Inches(0.3), width=Inches(1.0))
        
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.3))
    tf_f = footer_box.text_frame
    tf_f.word_wrap = True
    tf_f.margin_left = tf_f.margin_right = tf_f.margin_top = tf_f.margin_bottom = 0
    p_foot = tf_f.paragraphs[0]
    p_foot.text = "Edhigo Pani Technologies Private Limited  |  funding@dinasari.in  |  Confidential"
    p_foot.font.name = "Calibri"
    p_foot.font.size = Pt(9)
    p_foot.font.color.rgb = TEXT_MUTED

# Slide 1: Title Slide
def make_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_COLOR)
    
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(0.8), Inches(1.0), width=Inches(2.2))
        
    tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(5.8), Inches(3.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p1 = tf.paragraphs[0]
    p1.text = "Dinasari"
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(54)
    p1.font.bold = True
    p1.font.color.rgb = BRAND_GREEN
    
    p2 = tf.add_paragraph()
    p2.text = "Connecting Rural India's Workforce"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_DARK
    p2.space_before = Pt(8)
    p2.space_after = Pt(24)
    
    p3 = tf.add_paragraph()
    p3.text = "Digitizing Agricultural Labor & Machinery\nINVESTOR PITCH DECK | 2026"
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(14)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(12)
    
    img_path = os.path.join(ARTIFACT_DIR, "title_farmer_1783683567226.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(7.2), Inches(0.5), width=Inches(5.6), height=Inches(6.5))
        
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.3))
    tf_f = footer_box.text_frame
    tf_f.margin_left = tf_f.margin_right = tf_f.margin_top = tf_f.margin_bottom = 0
    p_foot = tf_f.paragraphs[0]
    p_foot.text = "Edhigo Pani Technologies Private Limited  |  funding@dinasari.in  |  Confidential"
    p_foot.font.name = "Calibri"
    p_foot.font.size = Pt(10)
    p_foot.font.color.rgb = TEXT_MUTED

# Slide 2: The Vision
def make_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide, "The Vision", "Our Vision: A Unified Digital Ecosystem for Rural India")
    
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(6.2), Inches(5.0))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p1 = tf.paragraphs[0]
    p1.text = "Edhigo Pani (Find Work)"
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = BRAND_ORANGE
    p1.space_after = Pt(12)
    
    p2 = tf.add_paragraph()
    p2.text = ("Dinasari is building a single, trusted platform to connect rural workers, farmers, "
               "machinery owners, and local services. We are digitizing the agricultural economy to "
               "create jobs, improve livelihoods, and strengthen village economies.")
    p2.font.name = "Calibri"
    p2.font.size = Pt(15)
    p2.font.color.rgb = TEXT_DARK
    p2.space_after = Pt(24)
    
    pillars = [
        ("Empower Workers", "Find reliable, well-paying work near you instantly.", BRAND_BLUE),
        ("Support Farmers", "Access labor and machinery on demand, eliminating seasonal stress.", BRAND_GREEN),
        ("Enable Growth", "Build a transparent, efficient, and digitized village economy.", BRAND_ORANGE)
    ]
    for title, desc, color in pillars:
        p = tf.add_paragraph()
        p.text = f"{title}: "
        p.font.name = "Segoe UI"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        
        run = p.add_run()
        run.text = desc
        run.font.name = "Calibri"
        run.font.size = Pt(14)
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(8)
        
    img_path = os.path.join(ARTIFACT_DIR, "vision_split_1783683581281.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(7.0), Inches(1.5), width=Inches(5.8), height=Inches(5.0))

# Slide 3: The Problem
def make_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide, "The Problem", "The Agricultural Labor Crisis: A Multi-Faceted Problem")
    
    left_margin = Inches(0.5)
    top_start = Inches(1.5)
    card_width = Inches(6.5)
    card_height = Inches(1.2)
    card_gap = Inches(0.2)
    
    problems = [
        ("Labor Shortages & Matchmaking", "Sowing & harvesting create massive labor demand, but finding workers relies on word-of-mouth or middleman contractors who take high commissions.", BRAND_RED),
        ("Underutilized Machinery & No Trust", "Tractor/harvester owners struggle to find bookings, while paper attendance tracking leads to disputes over hours worked and rates.", BRAND_RED),
        ("Payment Insecurity & Literacy Barriers", "Gig-workers suffer from delayed/defaulted wages, and digital platforms fail due to complex language and literacy barriers.", BRAND_RED)
    ]
    
    for idx, (title, desc, color) in enumerate(problems):
        y_pos = top_start + idx * (card_height + card_gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_margin, y_pos, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_COLOR
        card.line.color.rgb = LINE_COLOR
        card.line.width = Pt(1)
        
        indicator = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_margin + Inches(0.05), y_pos + Inches(0.05), Inches(0.1), card_height - Inches(0.1))
        indicator.fill.solid()
        indicator.fill.fore_color.rgb = color
        indicator.line.fill.background()
        
        tf = add_textbox(slide, left_margin + Inches(0.3), y_pos + Inches(0.1), card_width - Inches(0.5), card_height - Inches(0.2))
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_DARK
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Calibri"
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(4)
        
    stat_y = top_start + 3 * (card_height + card_gap)
    stat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_margin, stat_y, card_width, Inches(1.1))
    stat_box.fill.solid()
    stat_box.fill.fore_color.rgb = RGBColor(255, 240, 240)
    stat_box.line.color.rgb = BRAND_RED
    stat_box.line.width = Pt(1)
    
    tf_s = add_textbox(slide, left_margin + Inches(0.2), stat_y + Inches(0.1), card_width - Inches(0.4), Inches(0.9))
    p_s1 = tf_s.paragraphs[0]
    p_s1.text = "KEY STAT: Over 65% of India's population lives in rural areas, where access to organized employment remains limited."
    p_s1.font.name = "Segoe UI"
    p_s1.font.size = Pt(11.5)
    p_s1.font.bold = True
    p_s1.font.color.rgb = BRAND_RED
    
    p_s2 = tf_s.add_paragraph()
    p_s2.text = "(Data source: World Bank & Census of India 2011)"
    p_s2.font.name = "Calibri"
    p_s2.font.size = Pt(9.5)
    p_s2.font.color.rgb = TEXT_MUTED
    p_s2.space_before = Pt(4)
    
    img_path = os.path.join(ARTIFACT_DIR, "labor_crisis_1783683592623.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(7.3), Inches(1.5), width=Inches(5.5), height=Inches(5.0))

# Slide 4: The Solution
def make_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide, "The Solution", "Introducing Dinasari: The Unified Digital Ecosystem")
    
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.0))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p1 = tf.paragraphs[0]
    p1.text = "We are a one-stop solution that digitizes every step of the agricultural employment cycle."
    p1.font.name = "Calibri"
    p1.font.size = Pt(15)
    p1.font.color.rgb = TEXT_DARK
    p1.space_after = Pt(20)
    
    features = [
        ("Smart Matchmaking", "Farmers post jobs, and workers/groups discover nearby opportunities.", BRAND_GREEN),
        ("QR-Based Attendance", "GPS-verified QR check-in/out creates a tamper-proof digital ledger, eliminating disputes.", BRAND_BLUE),
        ("Direct Payouts", "Payments are calculated automatically and settled via cash or UPI, ensuring transparency.", BRAND_GREEN),
        ("Voice & Local Languages", "The app guides users in native languages (Telugu, Hindi, English) with simple iconography.", BRAND_ORANGE)
    ]
    
    for title, desc, color in features:
        p = tf.add_paragraph()
        p.text = f"• {title}: "
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        
        run = p.add_run()
        run.text = desc
        run.font.name = "Calibri"
        run.font.size = Pt(13)
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(10)
        
    p_stage = tf.add_paragraph()
    p_stage.text = "CURRENT STAGE: Product Ready & Pilot Ready"
    p_stage.font.name = "Segoe UI"
    p_stage.font.size = Pt(13.5)
    p_stage.font.bold = True
    p_stage.font.color.rgb = BRAND_GREEN
    p_stage.space_before = Pt(14)
    
    p_inn = tf.add_paragraph()
    p_inn.text = "INNOVATION: QR-based attendance, verified worker profiles, digital payments, and AI job matching."
    p_inn.font.name = "Segoe UI"
    p_inn.font.size = Pt(12)
    p_inn.font.bold = True
    p_inn.font.color.rgb = TEXT_DARK
    p_inn.space_before = Pt(4)
    
    img_path = os.path.join(ARTIFACT_DIR, "solution_ecosystem_1783683609054.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(7.3), Inches(1.5), width=Inches(5.5), height=Inches(5.0))

# Slide 5: Market Opportunity
def make_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide, "Market Opportunity", "A Massive, Untapped Rural Opportunity")
    
    # TAM Card
    card_tam = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(6.2), Inches(1.3))
    card_tam.fill.solid()
    card_tam.fill.fore_color.rgb = CARD_COLOR
    card_tam.line.color.rgb = BRAND_GREEN
    card_tam.line.width = Pt(2)
    
    tf_tam = add_textbox(slide, Inches(0.7), Inches(1.7), Inches(5.8), Inches(1.1))
    p1 = tf_tam.paragraphs[0]
    p1.text = "TAM (Total Addressable Market): ₹2.5 Lakh Cr+ ($30 Billion+)"
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = BRAND_GREEN
    p2 = tf_tam.add_paragraph()
    p2.text = "Total value of agricultural daily labor payrolls and equipment rentals across India."
    p2.font.name = "Calibri"
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_before = Pt(4)
    
    # SAM Card
    card_sam = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(3.1), Inches(5.8), Inches(1.3))
    card_sam.fill.solid()
    card_sam.fill.fore_color.rgb = CARD_COLOR
    card_sam.line.color.rgb = BRAND_BLUE
    card_sam.line.width = Pt(2)
    
    tf_sam = add_textbox(slide, Inches(1.1), Inches(3.2), Inches(5.4), Inches(1.1))
    p1_s = tf_sam.paragraphs[0]
    p1_s.text = "SAM (Serviceable Addressable Market): ₹25,000 Cr+ ($3 Billion+)"
    p1_s.font.name = "Segoe UI"
    p1_s.font.size = Pt(14)
    p1_s.font.bold = True
    p1_s.font.color.rgb = BRAND_BLUE
    p2_s = tf_sam.add_paragraph()
    p2_s.text = "Market in initial target states (AP, Telangana, Karnataka) with strong smartphone penetration and FPO networks."
    p2_s.font.name = "Calibri"
    p2_s.font.size = Pt(11)
    p2_s.font.color.rgb = TEXT_MUTED
    p2_s.space_before = Pt(4)
    
    # SOM Card
    card_som = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(4.6), Inches(5.4), Inches(1.3))
    card_som.fill.solid()
    card_som.fill.fore_color.rgb = CARD_COLOR
    card_som.line.color.rgb = BRAND_ORANGE
    card_som.line.width = Pt(2)
    
    tf_som = add_textbox(slide, Inches(1.5), Inches(4.7), Inches(5.0), Inches(1.1))
    p1_o = tf_som.paragraphs[0]
    p1_o.text = "SOM: ₹250 Cr+ ($30 Million+) in 3-5 Years"
    p1_o.font.name = "Segoe UI"
    p1_o.font.size = Pt(14)
    p1_o.font.bold = True
    p1_o.font.color.rgb = BRAND_ORANGE
    p2_o = tf_som.add_paragraph()
    p2_o.text = "Target transactional volume to be captured through focused operations, FPO partnerships, and viral growth."
    p2_o.font.name = "Calibri"
    p2_o.font.size = Pt(11)
    p2_o.font.color.rgb = TEXT_MUTED
    p2_o.space_before = Pt(4)
    
    # Right: Why Now?
    why_now = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.6), Inches(5.6), Inches(4.3))
    why_now.fill.solid()
    why_now.fill.fore_color.rgb = CARD_COLOR
    why_now.line.color.rgb = LINE_COLOR
    why_now.line.width = Pt(1)
    
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.2), Inches(1.6), Inches(5.6), Inches(0.4))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = BRAND_GREEN
    accent_bar.line.fill.background()
    
    tf_ab = add_textbox(slide, Inches(7.4), Inches(1.65), Inches(5.2), Inches(0.3))
    p_ab = tf_ab.paragraphs[0]
    p_ab.text = "WHY NOW?"
    p_ab.font.name = "Segoe UI"
    p_ab.font.size = Pt(12)
    p_ab.font.bold = True
    p_ab.font.color.rgb = TEXT_LIGHT
    
    tf_why = add_textbox(slide, Inches(7.4), Inches(2.2), Inches(5.2), Inches(3.5))
    p_why = tf_why.paragraphs[0]
    p_why.text = "India's Rural Digital Adoption"
    p_why.font.name = "Segoe UI"
    p_why.font.size = Pt(18)
    p_why.font.bold = True
    p_why.font.color.rgb = TEXT_DARK
    p_why.space_after = Pt(12)
    
    reasons = [
        ("Rapid Rural Digitization", "Internet growth in rural markets is now outpacing urban India, driven by widespread digital infrastructure."),
        ("Smartphone Adoption", "Increasing availability of low-cost smartphones has put app accessibility directly in the hands of farmers and laborers."),
        ("UPI Settlements", "Widespread adoption of the Unified Payments Interface (UPI) makes micro-transactions and wage splits secure and instant.")
    ]
    for title, text in reasons:
        p_res = tf_why.add_paragraph()
        p_res.text = f"• {title}: "
        p_res.font.name = "Segoe UI"
        p_res.font.size = Pt(12)
        p_res.font.bold = True
        p_res.font.color.rgb = BRAND_GREEN
        
        run = p_res.add_run()
        run.text = text
        run.font.name = "Calibri"
        run.font.size = Pt(12)
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
        p_res.space_after = Pt(8)

# Slide 6: User Journey - Farmer
def make_slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide, "User Journey - Farmer", "Farmer Journey: Sourcing Labor in Seconds")
    
    steps = [
        ("01", "Post Job Details", "Select crop, work type (Harvesting, Sowing), daily wage, and date."),
        ("02", "Match & Select", "Browse nearby individual workers or organized groups on our Live Discovery Map."),
        ("03", "QR Code Check-In", "Worker scans the farmer's QR code upon arrival, verifying location and start time."),
        ("04", "Track & Check-Out", "Monitor work status. Worker scans the QR code again at the end of the shift."),
        ("05", "Dynamic Settlement", "Review the automated payment summary, pay via Cash/UPI, and rate the worker.")
    ]
    
    card_width = Inches(2.2)
    card_height = Inches(3.8)
    card_gap = Inches(0.24)
    start_left = Inches(0.5)
    start_top = Inches(2.0)
    
    for idx, (num, title, desc) in enumerate(steps):
        left_pos = start_left + idx * (card_width + card_gap)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, start_top, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_COLOR
        card.line.color.rgb = BRAND_GREEN
        card.line.width = Pt(1.5)
        
        bubble = slide.shapes.add_shape(MSO_SHAPE.OVAL, left_pos + Inches(0.2), start_top - Inches(0.25), Inches(0.5), Inches(0.5))
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = BRAND_GREEN
        bubble.line.fill.background()
        
        tf_b = bubble.text_frame
        tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = 0
        p_b = tf_b.paragraphs[0]
        p_b.text = num
        p_b.font.name = "Segoe UI"
        p_b.font.size = Pt(14)
        p_b.font.bold = True
        p_b.font.color.rgb = TEXT_LIGHT
        p_b.alignment = PP_ALIGN.CENTER
        
        tf = add_textbox(slide, left_pos + Inches(0.15), start_top + Inches(0.4), card_width - Inches(0.3), card_height - Inches(0.5))
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_DARK
        p1.alignment = PP_ALIGN.CENTER
        p1.space_after = Pt(10)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Calibri"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED
        p2.alignment = PP_ALIGN.CENTER
        
        if idx < 4:
            arrow_left = left_pos + card_width + Inches(0.04)
            arrow_top = start_top + card_height / 2 - Inches(0.15)
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, arrow_left, arrow_top, Inches(0.16), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = LINE_COLOR
            arrow.line.fill.background()

# Slide 7: User Journey - Worker & Group Leader
def make_slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide, "User Journey - Worker & Leader", "Worker & Group Leader Journeys")
    
    # Left Column: Worker (Blue Theme)
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(5.9), Inches(5.0))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = CARD_COLOR
    left_card.line.color.rgb = BRAND_BLUE
    left_card.line.width = Pt(1.5)
    
    left_header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.6), Inches(5.9), Inches(0.5))
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = BRAND_BLUE
    left_header.line.fill.background()
    
    tf_lh = add_textbox(slide, Inches(0.7), Inches(1.65), Inches(5.5), Inches(0.4))
    p_lh = tf_lh.paragraphs[0]
    p_lh.text = "WORKER JOURNEY: RELIABLE WORK & PAYOUTS"
    p_lh.font.name = "Segoe UI"
    p_lh.font.size = Pt(11)
    p_lh.font.bold = True
    p_lh.font.color.rgb = TEXT_LIGHT
    
    tf_w = add_textbox(slide, Inches(0.7), Inches(2.2), Inches(5.5), Inches(4.3))
    worker_steps = [
        ("Profile Registration", "Register with mobile OTP, age, gender, and skills. Toggle 'Online' to appear in searches."),
        ("Push Notifications", "Receive loud, voice-guided job alerts detailing the farm distance, wage, and task."),
        ("GPS Field Navigation", "Accept the offer and view precise farm location routes."),
        ("QR Check-In/Out", "Scan the farmer's QR code to log attendance."),
        ("Earnings Ledger", "Track daily wages and settlement statuses on your personal dashboard.")
    ]
    for title, desc in worker_steps:
        p = tf_w.add_paragraph()
        p.text = f"• {title}: "
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = BRAND_BLUE
        
        run = p.add_run()
        run.text = desc
        run.font.name = "Calibri"
        run.font.size = Pt(11.5)
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(8)
        
    # Right Column: Group Leader (Orange Theme)
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.6), Inches(5.9), Inches(5.0))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = CARD_COLOR
    right_card.line.color.rgb = BRAND_ORANGE
    right_card.line.width = Pt(1.5)
    
    right_header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.6), Inches(5.9), Inches(0.5))
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = BRAND_ORANGE
    right_header.line.fill.background()
    
    tf_rh = add_textbox(slide, Inches(7.1), Inches(1.65), Inches(5.5), Inches(0.4))
    p_rh = tf_rh.paragraphs[0]
    p_rh.text = "GROUP LEADER JOURNEY: SCALING UP OPERATIONS"
    p_rh.font.name = "Segoe UI"
    p_rh.font.size = Pt(11)
    p_rh.font.bold = True
    p_rh.font.color.rgb = TEXT_LIGHT
    
    tf_l = add_textbox(slide, Inches(7.1), Inches(2.2), Inches(5.5), Inches(4.3))
    leader_steps = [
        ("Build Digital Groups", "Leaders (Mestris) register and bundle village workers into active groups."),
        ("Accept Collective Jobs", "Review and accept large-scale group jobs posted by farmers."),
        ("One-Scan Attendance", "Farmer scans the Leader's Group QR code, instantly checking in all team members."),
        ("Wage Splitter", "Generate structured salary breakdowns per worker based on individual attendance.")
    ]
    for title, desc in leader_steps:
        p = tf_l.add_paragraph()
        p.text = f"• {title}: "
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = BRAND_ORANGE
        
        run = p.add_run()
        run.text = desc
        run.font.name = "Calibri"
        run.font.size = Pt(11.5)
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(12)

# Slide 8: User Journey - Machinery Renter
def make_slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide, "User Journey - Machinery Renter", "Renters: Maximizing Tractor & Equipment Usage")
    
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.0))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p1 = tf.paragraphs[0]
    p1.text = "Dinasari is not just for labor; it's a full-fledged agricultural equipment marketplace."
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = BRAND_GREEN
    p1.space_after = Pt(16)
    
    features = [
        ("Equipment Listings", "Machinery owners upload details of tractors, harvesters, or power sprayers."),
        ("Set Pricing", "Control pricing by the hour or acre."),
        ("Real-Time Bookings", "Farmers discover and request machinery directly on the map."),
        ("QR Start/End Work", "QR scans verify the start and end of machinery rental operations."),
        ("Booking Calendar", "View upcoming bookings, block holiday dates, and manage logistics.")
    ]
    
    for title, desc in features:
        p = tf.add_paragraph()
        p.text = f"• {title}: "
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = BRAND_GREEN
        
        run = p.add_run()
        run.text = desc
        run.font.name = "Calibri"
        run.font.size = Pt(12.5)
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(8)
        
    img_path = os.path.join(ARTIFACT_DIR, "machinery_renter_1783683626452.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(7.3), Inches(1.5), width=Inches(5.5), height=Inches(5.0))

# Slide 9: Revenue Streams
def make_slide_9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide, "Revenue Streams", "Four Direct Revenue Streams")
    
    streams = [
        ("Platform Convenience Fee", "3% - 5%", "Transaction commission charged on daily worker payouts, calculated automatically on job completion.", BRAND_GREEN),
        ("Machinery Rental Commission", "8% - 10%", "Commission on heavy equipment rental bookings (tractors, tillers, harvest tools).", BRAND_BLUE),
        ("B2B Advertising Banners", "Promoted Ads", "Banner ads and promoted listings for seed brands, fertilizer suppliers, and tractor agencies.", BRAND_ORANGE),
        ("Financial Service Cross-Selling", "Value Addition", "Partnership margins for distributing micro-insurance (crop protection, accident) and agricultural credit.", TEXT_DARK)
    ]
    
    card_w = Inches(5.9)
    card_h = Inches(2.3)
    positions = [
        (Inches(0.5), Inches(1.6)),
        (Inches(6.9), Inches(1.6)),
        (Inches(0.5), Inches(4.3)),
        (Inches(6.9), Inches(4.3))
    ]
    
    for idx, (title, rate, desc, color) in enumerate(streams):
        left_pos, top_pos = positions[idx]
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_COLOR
        card.line.color.rgb = LINE_COLOR
        card.line.width = Pt(1)
        
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, top_pos, card_w, Inches(0.15))
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()
        
        tf = add_textbox(slide, left_pos + Inches(0.3), top_pos + Inches(0.3), card_w - Inches(0.6), card_h - Inches(0.4))
        
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Segoe UI"
        p_title.font.size = Pt(14)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_DARK
        
        p_rate = tf.add_paragraph()
        p_rate.text = rate
        p_rate.font.name = "Segoe UI"
        p_rate.font.size = Pt(28)
        p_rate.font.bold = True
        p_rate.font.color.rgb = color
        p_rate.space_before = Pt(4)
        p_rate.space_after = Pt(4)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Calibri"
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_MUTED

# Slide 10: Go-To-Market & Growth Strategy
def make_slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide, "Go-To-Market Strategy", "Localized GTM & Growth Strategy")
    
    strategies = [
        ("FPO Integration", "Farmer Cooperatives", "Partnering with Farmer Producer Organizations to register clusters of farmers efficiently.", BRAND_GREEN),
        ("Gram Panchayat campaigns", "Village 'Choupals'", "Conducting physical village meetings to build trust and address digital barriers.", BRAND_BLUE),
        ("Leader Referral Loop", "Mestri Virality", "Incentivizing gang leaders (Mestris) to onboard their regular worker teams, creating rapid growth.", BRAND_ORANGE),
        ("Retailer Affiliates", "Agri-Input Stores", "Onboarding local seed & fertilizer retailers as Dinasari registration kiosks.", TEXT_DARK)
    ]
    
    card_w = Inches(2.75)
    card_h = Inches(4.3)
    card_gap = Inches(0.3)
    start_left = Inches(0.5)
    start_top = Inches(1.8)
    
    for idx, (title, sub, desc, color) in enumerate(strategies):
        left_pos = start_left + idx * (card_w + card_gap)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, start_top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_COLOR
        card.line.color.rgb = color
        card.line.width = Pt(1.5)
        
        tf = add_textbox(slide, left_pos + Inches(0.15), start_top + Inches(0.2), card_w - Inches(0.3), card_h - Inches(0.4))
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = color
        p1.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = sub.upper()
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(9.5)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_MUTED
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(4)
        p2.space_after = Pt(14)
        
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.name = "Calibri"
        p3.font.size = Pt(11)
        p3.font.color.rgb = TEXT_MUTED
        p3.alignment = PP_ALIGN.CENTER
        
    bot_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.6))
    bot_box.fill.solid()
    bot_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
    bot_box.line.color.rgb = BRAND_BLUE
    bot_box.line.width = Pt(1)
    
    tf_b = add_textbox(slide, Inches(0.7), Inches(6.25), Inches(12.0), Inches(0.5))
    p_b = tf_b.paragraphs[0]
    p_b.text = "VIRAL CHANNELS: Localized SMS & Voice Broadcasts explain app benefits in native Telugu and Hindi directly to users."
    p_b.font.name = "Segoe UI"
    p_b.font.size = Pt(11.5)
    p_b.font.bold = True
    p_b.font.color.rgb = BRAND_BLUE

# Slide 11: Unit Economics & Projections
def make_slide_11(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide, "Unit Economics & Projections", "Strong Unit Economics & Rapid Scaling")
    
    economics = [
        ("CAC (Customer Acquisition Cost)", "₹200 ($2.4)", "Low CAC achieved via viral worker-to-worker loops and FPO channels.", BRAND_GREEN),
        ("LTV (3-Year Lifetime Value)", "₹1,800 ($22)", "Consistent transactional commissions across multiple crop seasons.", BRAND_BLUE),
        ("LTV / CAC Ratio", "9.0x", "High profitability and marketing efficiency.", BRAND_ORANGE)
    ]
    
    left_start = Inches(0.5)
    top_start = Inches(1.6)
    card_w = Inches(4.5)
    card_h = Inches(1.4)
    card_gap = Inches(0.2)
    
    for idx, (title, val, desc, color) in enumerate(economics):
        y_pos = top_start + idx * (card_h + card_gap)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_start, y_pos, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_COLOR
        card.line.color.rgb = LINE_COLOR
        card.line.width = Pt(1)
        
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_start, y_pos, Inches(0.12), card_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        tf = add_textbox(slide, left_start + Inches(0.3), y_pos + Inches(0.1), card_w - Inches(0.4), card_h - Inches(0.2))
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Segoe UI"
        p_title.font.size = Pt(11)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_DARK
        
        p_val = tf.add_paragraph()
        p_val.text = val
        p_val.font.name = "Segoe UI"
        p_val.font.size = Pt(22)
        p_val.font.bold = True
        p_val.font.color.rgb = color
        p_val.space_before = Pt(2)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Calibri"
        p_desc.font.size = Pt(9.5)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_before = Pt(2)
        
    rows, cols = 4, 4
    t_left = Inches(5.4)
    t_top = Inches(1.8)
    t_width = Inches(7.4)
    t_height = Inches(3.5)
    
    table_shape = slide.shapes.add_table(rows, cols, t_left, t_top, t_width, t_height)
    table = table_shape.table
    
    table.columns[0].width = Inches(2.6)
    table.columns[1].width = Inches(1.6)
    table.columns[2].width = Inches(1.6)
    table.columns[3].width = Inches(1.6)
    
    headers = ["Metric", "Year 1", "Year 2", "Year 3"]
    data = [
        ["Active Users", "50,000", "250,000", "1,000,000+"],
        ["Gross Transaction Value (GTV)", "₹50 Cr", "₹300 Cr", "₹1,400 Cr"],
        ["Platform Revenue", "₹2.5 Cr", "₹15 Cr", "₹70 Cr"]
    ]
    
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_GREEN
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER
        
    for row_idx, row_data in enumerate(data):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            else:
                cell.fill.fore_color.rgb = RGBColor(240, 245, 240)
            
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = "Calibri"
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_DARK
            if col_idx == 0:
                p.font.bold = True
                p.alignment = PP_ALIGN.LEFT
                p.font.name = "Segoe UI"
            else:
                p.alignment = PP_ALIGN.CENTER
                if col_idx == 3:
                    p.font.bold = True
                    p.font.color.rgb = BRAND_GREEN

# Slide 12: Competitive Moat
def make_slide_12(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide, "Competitive Moat", "Strong Competitive Moat in Rural Tech")
    
    rows, cols = 7, 4
    left = Inches(1.2)
    top = Inches(1.6)
    width = Inches(10.9)
    height = Inches(4.8)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    table.columns[0].width = Inches(4.3)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(2.2)
    table.columns[3].width = Inches(2.2)
    
    headers = ["Feature / Advantage", "Dinasari", "Job Portals", "Contractors"]
    data = [
        ["Rural Agri-Specific Workflows", "✅ Yes", "❌ No (Urban Focus)", "✅ Yes (Inefficient)"],
        ["Group-Hiring Mechanics", "✅ Yes", "❌ No (1-to-1)", "✅ Yes (Exploitative)"],
        ["QR-Based Attendance Ledger", "✅ Yes", "❌ No", "❌ No"],
        ["Machinery Rental Integration", "✅ Yes", "❌ No", "❌ No"],
        ["Grievance Redressal Desk", "✅ Yes", "❌ No", "❌ No"],
        ["Low-Literacy Voice & Icon Design", "✅ Yes", "❌ No", "❌ No"]
    ]
    
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        if col_idx == 1:
            cell.fill.fore_color.rgb = BRAND_GREEN
        else:
            cell.fill.fore_color.rgb = TEXT_DARK
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER
        
    for row_idx, row_data in enumerate(data):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            if col_idx == 1:
                cell.fill.fore_color.rgb = RGBColor(240, 248, 240)
            elif row_idx % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            else:
                cell.fill.fore_color.rgb = RGBColor(245, 245, 245)
                
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = "Calibri"
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_DARK
            
            if col_idx == 0:
                p.alignment = PP_ALIGN.LEFT
                p.font.bold = True
                p.font.name = "Segoe UI"
            else:
                p.alignment = PP_ALIGN.CENTER
                if col_idx == 1:
                    p.font.bold = True
                    p.font.color.rgb = BRAND_GREEN
                elif "✅" in val:
                    p.font.color.rgb = BRAND_GREEN
                elif "❌" in val:
                    p.font.color.rgb = BRAND_RED

# Slide 13: The Team
def make_slide_13(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide, "The Team", "The Team Behind Dinasari")
    
    team = [
        ("Revanth Lokesh", "CEO & Co-founder", "Driving the mission to empower rural India through technology.", "revanth_avatar_1783683643803.png", BRAND_GREEN),
        ("Vishnu Sai", "CTO & Co-founder", "Responsible for platform technology, product innovation, and digital infrastructure.", "vishnu_avatar_1783683656919.png", BRAND_BLUE),
        ("Aruna Mani", "Mentor & Advisor", "Manager at SmartFX (Fintech), providing strategic guidance on business growth, operations, and execution.", "aruna_avatar_1783683668917.png", BRAND_ORANGE)
    ]
    
    card_w = Inches(3.6)
    card_h = Inches(5.0)
    card_gap = Inches(0.5)
    start_left = Inches(0.5)
    start_top = Inches(1.6)
    
    for idx, (name, role, desc, img_name, color) in enumerate(team):
        left_pos = start_left + idx * (card_w + card_gap)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, start_top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_COLOR
        card.line.color.rgb = LINE_COLOR
        card.line.width = Pt(1)
        
        img_path = os.path.join(ARTIFACT_DIR, img_name)
        img_size = Inches(1.8)
        img_left = left_pos + (card_w - img_size) / 2
        img_top = start_top + Inches(0.4)
        
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, img_left, img_top, width=img_size, height=img_size)
            ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, img_left, img_top, img_size, img_size)
            ring.fill.background()
            ring.line.color.rgb = color
            ring.line.width = Pt(2.5)
            
        tf = add_textbox(slide, left_pos + Inches(0.2), start_top + Inches(2.4), card_w - Inches(0.4), Inches(2.4))
        
        p_name = tf.paragraphs[0]
        p_name.text = name
        p_name.font.name = "Segoe UI"
        p_name.font.size = Pt(18)
        p_name.font.bold = True
        p_name.font.color.rgb = TEXT_DARK
        p_name.alignment = PP_ALIGN.CENTER
        
        p_role = tf.add_paragraph()
        p_role.text = role.upper()
        p_role.font.name = "Segoe UI"
        p_role.font.size = Pt(11)
        p_role.font.bold = True
        p_role.font.color.rgb = color
        p_role.alignment = PP_ALIGN.CENTER
        p_role.space_before = Pt(4)
        p_role.space_after = Pt(12)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Calibri"
        p_desc.font.size = Pt(11.5)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.alignment = PP_ALIGN.CENTER

# Slide 14: Key Risks & Mitigation
def make_slide_14(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide, "Risks & Mitigation", "Key Risks & Mitigation Strategy")
    
    rows, cols = 5, 2
    left = Inches(1.2)
    top = Inches(1.6)
    width = Inches(10.9)
    height = Inches(4.8)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    table.columns[0].width = Inches(3.9)
    table.columns[1].width = Inches(7.0)
    
    headers = ["Risk", "Mitigation Strategy"]
    data = [
        ["Slow Digital Adoption", "Drive adoption through village ambassadors, FPOs, and local awareness campaigns (Choupals)."],
        ["Building Trust", "Use verified profiles, QR-based attendance, ratings, and transparent, instant payments."],
        ["Scaling Operations", "Expand district by district through strategic partnerships and a scalable cloud-based platform."],
        ["Maintaining Quality", "Implement a dedicated dispute resolution desk and in-field coordinators."]
    ]
    
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_GREEN
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_LIGHT
        p.alignment = PP_ALIGN.LEFT
        
    for row_idx, row_data in enumerate(data):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            else:
                cell.fill.fore_color.rgb = RGBColor(245, 245, 245)
                
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = "Calibri"
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_DARK
            
            if col_idx == 0:
                p.font.bold = True
                p.font.name = "Segoe UI"
                p.font.color.rgb = BRAND_RED

# Slide 15: The Ask & Use of Funds
def make_slide_15(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide, "The Ask", "The Ask: Seed Funding")
    
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5.5), Inches(5.0))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p1 = tf.paragraphs[0]
    p1.text = "We are seeking ₹8 Crore ($1.0 Million USD) to capitalize on our growth."
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_DARK
    p1.space_after = Pt(24)
    
    p3 = tf.add_paragraph()
    p3.text = "3-Year Revenue Snapshot:"
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(16)
    p3.font.bold = True
    p3.font.color.rgb = BRAND_GREEN
    p3.space_after = Pt(12)
    
    revenues = [
        ("Year 1", "₹2.5 Cr"),
        ("Year 2", "₹15 Cr"),
        ("Year 3", "₹70 Cr")
    ]
    for yr, rev in revenues:
        p_yr = tf.add_paragraph()
        p_yr.text = f"• {yr}: "
        p_yr.font.name = "Segoe UI"
        p_yr.font.size = Pt(13)
        p_yr.font.bold = True
        p_yr.font.color.rgb = TEXT_DARK
        
        run = p_yr.add_run()
        run.text = rev
        run.font.name = "Calibri"
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = BRAND_GREEN
        p_yr.space_after = Pt(8)
        
    funds = [
        ("40% Product & Engineering", "Building scalable backend APIs, chatbot voice capabilities, and offline sync.", BRAND_GREEN),
        ("35% GTM & Customer Acquisition", "Marketing roadshows, village panchayat setups, and FPO incentive schemes.", BRAND_BLUE),
        ("15% Operations & Support", "Grievance Desk expansion, on-field coordinators, and helpline operators.", BRAND_ORANGE),
        ("10% General & Administrative", "Corporate legal compliance, office space, and reserve contingency.", TEXT_DARK)
    ]
    
    r_start = Inches(6.5)
    r_top = Inches(1.6)
    card_w = Inches(6.333)
    card_h = Inches(1.15)
    card_gap = Inches(0.12)
    
    for idx, (title, desc, color) in enumerate(funds):
        y_pos = r_top + idx * (card_h + card_gap)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, r_start, y_pos, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_COLOR
        card.line.color.rgb = LINE_COLOR
        card.line.width = Pt(1)
        
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, r_start, y_pos, Inches(0.15), card_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        tf_f = add_textbox(slide, r_start + Inches(0.3), y_pos + Inches(0.1), card_w - Inches(0.4), card_h - Inches(0.2))
        p_title = tf_f.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Segoe UI"
        p_title.font.size = Pt(12)
        p_title.font.bold = True
        p_title.font.color.rgb = color
        
        p_desc = tf_f.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Calibri"
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_before = Pt(2)

# Slide 16: SDG Alignment & Closing
def make_slide_16(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide, "Impact & Closing", "Our Impact: Aligned with National Priorities")
    
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(6.0), Inches(5.0))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p1 = tf.paragraphs[0]
    p1.text = "SDG & National Alignment"
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = BRAND_GREEN
    p1.space_after = Pt(12)
    
    sdgs = [
        ("Create Jobs", "Sustainable employment for millions of rural workers (SDG 1 & 8)."),
        ("Boost Agriculture", "Improve farmer productivity through faster workforce and machinery access (SDG 2 & 9)."),
        ("Promote Financial Inclusion", "Enable secure digital payments and credit access (SDG 8 & 10)."),
        ("Build Transparency", "Increase trust through verified profiles and secure transactions.")
    ]
    for title, desc in sdgs:
        p = tf.add_paragraph()
        p.text = f"• {title}: "
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = BRAND_GREEN
        
        run = p.add_run()
        run.text = desc
        run.font.name = "Calibri"
        run.font.size = Pt(11.5)
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(8)
        
    p_vande = tf.add_paragraph()
    p_vande.text = "Vande Bharatam Alignment: Supports Digital India, financial inclusion, employment generation, and rural development."
    p_vande.font.name = "Segoe UI"
    p_vande.font.size = Pt(11.5)
    p_vande.font.bold = True
    p_vande.font.color.rgb = TEXT_DARK
    p_vande.space_before = Pt(12)
    
    r_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.6), Inches(5.8), Inches(4.8))
    r_card.fill.solid()
    r_card.fill.fore_color.rgb = CARD_COLOR
    r_card.line.color.rgb = BRAND_GREEN
    r_card.line.width = Pt(2)
    
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(8.9), Inches(2.0), width=Inches(2.0))
        
    tf_c = add_textbox(slide, Inches(7.2), Inches(3.5), Inches(5.4), Inches(2.5))
    p_c1 = tf_c.paragraphs[0]
    p_c1.text = "Join us in building the future of rural India."
    p_c1.font.name = "Segoe UI"
    p_c1.font.size = Pt(18)
    p_c1.font.bold = True
    p_c1.font.color.rgb = TEXT_DARK
    p_c1.alignment = PP_ALIGN.CENTER
    p_c1.space_after = Pt(16)
    
    p_c2 = tf_c.add_paragraph()
    p_c2.text = "Contact:"
    p_c2.font.name = "Segoe UI"
    p_c2.font.size = Pt(11)
    p_c2.font.bold = True
    p_c2.font.color.rgb = TEXT_MUTED
    p_c2.alignment = PP_ALIGN.CENTER
    
    p_c3 = tf_c.add_paragraph()
    p_c3.text = "funding@dinasari.in"
    p_c3.font.name = "Segoe UI"
    p_c3.font.size = Pt(16)
    p_c3.font.bold = True
    p_c3.font.color.rgb = BRAND_GREEN
    p_c3.alignment = PP_ALIGN.CENTER
    p_c3.space_before = Pt(4)
    p_c3.space_after = Pt(8)
    
    p_c4 = tf_c.add_paragraph()
    p_c4.text = "Edhigo Pani Technologies Private Limited"
    p_c4.font.name = "Segoe UI"
    p_c4.font.size = Pt(12)
    p_c4.font.color.rgb = TEXT_DARK
    p_c4.alignment = PP_ALIGN.CENTER

# Build presentation
make_slide_1(prs)
make_slide_2(prs)
make_slide_3(prs)
make_slide_4(prs)
make_slide_5(prs)
make_slide_6(prs)
make_slide_7(prs)
make_slide_8(prs)
make_slide_9(prs)
make_slide_10(prs)
make_slide_11(prs)
make_slide_12(prs)
make_slide_13(prs)
make_slide_14(prs)
make_slide_15(prs)
make_slide_16(prs)

# Save
prs.save(OUTPUT_PATH)
print("Pitch deck successfully generated at:", OUTPUT_PATH)
